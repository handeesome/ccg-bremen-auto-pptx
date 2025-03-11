from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.dml.color import RGBColor
from pptx.util import Cm, Pt
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.text import PP_ALIGN
import re
import os


def addTextBox(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    # start a new line if sentence too long and it hits the boundrary
    txBox.text_frame.word_wrap = True
    return txBox


def newCenterBox(slide, content, fontSize=66):
    txBox = addTextBox(slide, 0, 3.21, 25.4, 12.62)
    p = txBox.text_frame.paragraphs[0]
    p.text = content
    setFont(p, fontSize)
    p.alignment = PP_ALIGN.CENTER
    # align the text to the middle vertically
    txBox.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return txBox


def setFont(p, size, fontName='黑体', isBlue=False):
    """
    Change the font styles for a paragraph.

    Parameters:
        p (paragraph): The paragraph whose font styles will be changed.
        size (int): Font size.
        fontName (str, optional): Font family. Default is '黑体'.
        isBlue (bool, optional): Set text to blue if True. Default is False.

    Returns:
        font: The modified font object.
    """
    font = p.font
    font.language_id = MSO_LANGUAGE_ID.SIMPLIFIED_CHINESE
    font.name = fontName
    font.size = Pt(size)
    if isBlue == True:
        font.color.rgb = RGBColor(0x2E, 0x75, 0xB6)

    return font


def newSlide(prs, title, fontSize=42):
    '''
    Create a new slide in a PowerPoint presentation with a title.

    Parameters:
        prs (Presentation): The PowerPoint presentation to which the slide will be added.
        title (str): The title to be displayed on the slide.
        fontSize (int, optional): The font size for the title. Default is 42.

    Returns:
        Slide: The newly created slide.

    Explanation:
        This function creates a new slide in the provided PowerPoint presentation (prs) with a specified title.
        The title is formatted with the given font size (fontSize) and is centered vertically.
    '''
    empty_slide = prs.slide_layouts[6]
    slide = prs.slides.add_slide(empty_slide)
    txBox = addTextBox(slide, 2.91, 0, 22.49, 3.08)
    p = txBox.text_frame.paragraphs[0]
    p.text = title
    setFont(p, fontSize, isBlue=True)
    txBox.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return slide


def verseGroup(verses, character_limit=90, isQiYing=False):
    """
    Group verses into segments, considering character limits for each group.

    Parameters:
        verses (list): List of verses to be grouped.
        character_limit (int, optional): The character limit for each group. Default is 90.
        isQiYing (bool, optional): 启应经文 if True. Default is False.

    Returns:
        list: A list of grouped verses, where each group respects the character limit.
    """
    groups = []
    current_group = []
    current_character_count = 0

    for verse in verses:
        verse_character_count = len(verse)

        if isQiYing and len(current_group) == 2:
            # If creating 启应经文 and current group has 2 verses, start a new group
            groups.append(current_group)
            current_group = [verse]
            current_character_count = verse_character_count
        elif not current_group:  # Ensure first group is NOT empty
            # Start the first group directly
            current_group.append(verse)
            current_character_count = verse_character_count
        elif current_character_count + verse_character_count <= character_limit:
            # If adding the current verse doesn't exceed the character limit, add it
            current_group.append(verse)
            current_character_count += verse_character_count
        else:
            # Start a new group as the current verse exceeds the character limit
            groups.append(current_group)
            current_group = [verse]
            current_character_count = verse_character_count

    # Add the last group if it's not empty
    if current_group:
        groups.append(current_group)

    return groups


def parseVerses(verses, isQiYing=False):
    verseArr = []
    verseFrom = int(verses["verseFrom"])
    for verse in verses["fullVerse"]:
        if isQiYing:# 启应经文 needs the verses without verse numbers
            verseArr.append(verse[1])
        else:
            verseArr.append(str(verseFrom) +'.'+verse[1])
            verseFrom += 1
    return verseArr

def divide_verse_evenly(verse):
    """
    Divide a verse into two chunks, attempting to split it evenly near the middle while considering delimiters.
    Used when a verse is too long in 启应经文, also when the length of verses is odd number, 
    there is need to split one verse into two.

    Parameters:
        verse (str): The verse to be divided.

    Returns:
        list: A list of two chunks after division.
    """
    delimiters = ['，', '。', '：', '；', '?', '！', '、']
    chunks = []
    middle = len(verse) // 2
    left_delimiter = None
    right_delimiter = None

    # Search for delimiter positions from the middle towards the left
    for i in range(middle, -1, -1):
        if verse[i] in delimiters:
            left_delimiter = i
            break
    # Search for delimiter positions from the middle towards the right
    for i in range(middle, len(verse)):
        if verse[i] in delimiters:
            right_delimiter = i
            break
    if left_delimiter is not None and right_delimiter is not None:
        if abs(middle-left_delimiter) <= abs(right_delimiter-middle):
            delimiter_index = left_delimiter
        else:
            delimiter_index = right_delimiter
    else:
        delimiter_index = right_delimiter
    chunks.append(verse[:delimiter_index + 1])
    chunks.append(verse[delimiter_index + 1:])
    return chunks


def addAudio(slide, audioPath):
    slide.shapes.add_movie(audioPath, Cm(20.14), Cm(16.79), Cm(
        5), Cm(2.28), poster_frame_image=None)


def lyricsPages(prs, song, lyrics, audioPath=None, isOnePage=False, titleFontSize=42):
    """
    Create PowerPoint slides for song lyrics and audio.

    Parameters:
        prs (Presentation): The PowerPoint presentation to which slides will be added.
        song (str): The title of the song.

        lyrics (str): The lyrics of the song, including page markers.
        i.e. 
        ""
        [Page 1]
        第一句歌词
        第二句歌词

        [Page 2]
        第二页歌词
        ...
        ""

        audioPath (str): The path to the audio file.
        isOnePage (bool, optional): lyrics is a single page if True. Default is False.
        titleFontSize (int, optional): The font size for slide titles. Default is 42.

    Returns:
        None
    """
    lyrics_pages = re.split(r'\[Page \d+\]', lyrics)
    lyrics_pages = [page.strip() for page in lyrics_pages if page.strip()]
    pageNum = len(lyrics_pages)
    for i in range(pageNum):
        title = song + f' {i+1}/{pageNum}'
        slide = newSlide(prs, title, titleFontSize)
        txBox = addTextBox(slide, 0, 3.68, 25.4, 12.09)
        p = txBox.text_frame.paragraphs[0]
        p.text = lyrics_pages[i]
        p.alignment = PP_ALIGN.CENTER
        setFont(p, 44)
    slide = prs.slides[-pageNum]
    if isOnePage:
        title = song
        p = slide.shapes[0].text_frame.paragraphs[0]
        p.text = title
        p = slide.shapes[1].text_frame.paragraphs[0]
        setFont(p, 48)
    if(audioPath):
        addAudio(slide, audioPath)


def itemsPagesHasSpace(prs, title, items, charPerLine, fontSize=36, isNumber=False, character_limit=90):
    """
    Creates PowerPoint slides with formatted text for a list of items with line breaks.

    Parameters:
        prs (Presentation): The PowerPoint presentation object.
        title (str): The title for the slides.
        items (list): List of items to be displayed on slides.
        charPerLine (int): Maximum number of characters per line.
        fontSize (int, optional): Font size for the text. Defaults to 36.
        isNumber (bool, optional): If True, adds item numbers. Defaults to False.

    Additional Info:
        This function adds extra space before every non-numbered line to ensure consistent
        vertical alignment. However, due to certain limitations, the implementation may not
        achieve perfect alignment in all cases.

    Note:
        The original version of this function is itemsPages. In cases where the extra space
        feature doesn't meet expectations, consider using itemsPages for the standard layout.
    """
    number = 0

    groups = verseGroup(items, character_limit=character_limit)
    for group in groups:
        slide = newSlide(prs, title)
        txBox = addTextBox(slide, 1.72, 3.45, 21.8, 15.61)
        txBox.text_frame.clear()

        newLine = '\n  '
        for item in group:
            p = txBox.text_frame.add_paragraph()
            if isNumber:
                # number the items
                number += 1
                numbered_item = f'{number} ' + item
            numbered_item = item
            pattern = re.compile(r'(\d+).(.*)')
            match = re.search(pattern, numbered_item)
            if match:
                number = match.group(1)
                if int(number)>9:
                    newLine = '\n   '
                text_part = match.group(2)
                delimiters = ['，', '。', '：', '；', '?', '、', '！']
                substrings = []
                i = 0
                while i < len(text_part):
                    line = text_part[i:i+charPerLine]
                    if substrings and line[0] in delimiters:
                        # If the first character of the line is a delimiter, append it to the previous line
                        substrings[-1] += line[0]
                        i += 1  # Move to the next character
                    else:
                        # Otherwise, add the line as a new entry
                        substrings.append(line)
                        i += charPerLine  # Move to the next line
                p.text = number + '.' + newLine.join(substrings)
                setFont(p, fontSize)


def itemsPages(prs, title, items, fontSize, isNumber=False, character_limit=90):
    """
    Creates PowerPoint slides with formatted text for a list of items.

    Parameters:
        prs (Presentation): The PowerPoint presentation object.
        title (str): The title for the slides.
        items (list): List of items to be displayed on slides.
        fontSize (int): Font size for the text.
        isNumber (bool, optional): If True, adds item numbers. Defaults to False.
        character_limit (int, optional): Maximum character limit for each group. Defaults to 90.
    """
    # add numbers to every item if needed
    if isNumber:
        modified_items = []
        for i, item in enumerate(items, 1):
            modified_item = f'{i}. {item}'
            modified_items.append(modified_item)

        groups = verseGroup(modified_items, character_limit=character_limit)
    else:
        groups = verseGroup(items, character_limit=character_limit)
    for group in groups:
        slide = newSlide(prs, title)
        txBox = addTextBox(slide, 1.72, 3.45, 21.8, 15.61)
        txBox.text_frame.clear()

        for item in group:
            p = txBox.text_frame.add_paragraph()
            p.text = item
            setFont(p, fontSize)
            p.space_after = Pt(10)