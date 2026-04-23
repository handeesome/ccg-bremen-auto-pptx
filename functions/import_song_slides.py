import os
import re
import shutil
import tempfile
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path

from pptx import Presentation

PPTX_EXTENSION = ".pptx"
SLIDE_TITLE_PATTERN = re.compile(r"^.+\s+\d+/\d+\s*$")
AUDIO_EXTENSIONS = {".mp3", ".wav", ".m4a", ".aac", ".wma", ".ogg", ".mid", ".midi"}
REL_NS = {"rel": "http://schemas.openxmlformats.org/package/2006/relationships"}


def sanitize_song_name(filename):
    stem = Path(filename).stem.strip() or "song"
    sanitized = re.sub(r'[<>:"/\\\\|?*]+', "_", stem)
    sanitized = re.sub(r"\s+", " ", sanitized).strip(" .")
    return sanitized or "song"


def extract_song_slides(uploaded_file):
    original_name = uploaded_file.filename or "song"
    with tempfile.TemporaryDirectory() as temp_dir:
        source_path = os.path.join(temp_dir, original_name)
        uploaded_file.save(source_path)
        return extract_song_slides_from_path(source_path, original_name)


def extract_song_slides_from_path(source_path, original_name):
    suffix = Path(original_name).suffix.lower()
    if suffix != PPTX_EXTENSION:
        raise ValueError("Only .pptx files are supported.")

    slides, detected_title = read_pptx_slide_texts(source_path)
    return {
        "originalName": original_name,
        "suggestedTitle": sanitize_song_name(detected_title or original_name),
        "slides": slides,
    }


def read_pptx_slide_texts(file_path):
    presentation = Presentation(file_path)
    slides = []
    detected_title = None

    for index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            texts.extend(extract_text_from_shape(shape))

        raw_slide_text = "\n".join(line for line in texts if line)
        if detected_title is None:
            detected_title = extract_song_title(raw_slide_text)
        slide_text = normalize_slide_text(raw_slide_text)
        slides.append(
            {
                "number": index,
                "text": slide_text,
                "lines": slide_text.splitlines() if slide_text else [],
            }
        )

    return slides, detected_title


def extract_text_from_shape(shape):
    texts = []

    if getattr(shape, "has_text_frame", False):
        paragraphs = []
        for paragraph in shape.text_frame.paragraphs:
            line = normalize_line_breaks("".join(run.text for run in paragraph.runs)).strip()
            if not line:
                line = normalize_line_breaks(paragraph.text).strip()
            if line:
                paragraphs.append(line)
        texts.extend(paragraphs)

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            row_values = []
            for cell in row.cells:
                value = cell.text.strip()
                if value:
                    row_values.append(value)
            if row_values:
                texts.append(" | ".join(row_values))

    if shape.shape_type == 6:
        for child in shape.shapes:
            texts.extend(extract_text_from_shape(child))

    return texts


def normalize_line_breaks(value):
    return str(value).replace("\r\n", "\n").replace("\r", "\n").replace("\x0b", "\n")


def normalize_slide_text(value):
    normalized = normalize_line_breaks(value)
    lines = [line.strip() for line in normalized.split("\n")]

    while lines and not lines[0]:
        lines.pop(0)
    while lines and not lines[-1]:
        lines.pop()

    if lines and SLIDE_TITLE_PATTERN.match(lines[0]):
        lines = lines[1:]

    cleaned_lines = []
    previous_blank = False
    for line in lines:
        is_blank = line == ""
        if is_blank and previous_blank:
            continue
        cleaned_lines.append(line)
        previous_blank = is_blank

    return "\n".join(cleaned_lines).strip()


def extract_song_title(value):
    normalized = normalize_line_breaks(value)
    for line in normalized.split("\n"):
        candidate = line.strip()
        if not candidate:
            continue
        if SLIDE_TITLE_PATTERN.match(candidate):
            return re.sub(r"\s+\d+/\d+\s*$", "", candidate).strip()
    return None


def cleanup_paths(paths):
    for path in paths:
        if not path:
            continue
        if os.path.isdir(path):
            shutil.rmtree(path, ignore_errors=True)
        elif os.path.exists(path):
            os.remove(path)


def extract_slide_one_audio(pptx_path, temp_dir):
    rels_path = "ppt/slides/_rels/slide1.xml.rels"
    try:
        with zipfile.ZipFile(pptx_path, "r") as archive:
            if rels_path not in archive.namelist():
                return None

            rels_root = ET.fromstring(archive.read(rels_path))
            for relationship in rels_root.findall("rel:Relationship", REL_NS):
                target = relationship.attrib.get("Target", "")
                if not target.startswith("../media/"):
                    continue

                media_path = target.replace("../", "ppt/", 1)
                extension = Path(media_path).suffix.lower()
                if extension not in AUDIO_EXTENSIONS:
                    continue
                if media_path not in archive.namelist():
                    continue

                output_name = f"{Path(pptx_path).stem}_slide1_audio{extension}"
                output_path = os.path.join(temp_dir, output_name)
                with open(output_path, "wb") as audio_file:
                    audio_file.write(archive.read(media_path))
                return output_path
    except (zipfile.BadZipFile, ET.ParseError):
        return None

    return None
