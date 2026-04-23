from pptx import Presentation
from .helper import *

def generate_lyrics_pptx(templatePath, pages, songName, output_dir=None, audio_path=None):
    title = songName
    prs = Presentation(templatePath)

    # insert [page 1] and [page 2] to the beginning of pages in loop
    pages_with_markers = []
    for i, page in enumerate(pages):
        pages_with_markers.append(f"[Page {i+1}]\n{page}")
    # concatenate pages with pages
    pages = "\n".join(pages_with_markers)

    lyricsPages(prs, title, pages)
    if audio_path and len(prs.slides) > 0:
        addAudio(prs.slides[0], audio_path)
    fileName = f"{title}.pptx"
    output_dir = output_dir or os.path.join("static", "temp")
    os.makedirs(output_dir, exist_ok=True)
    prs.save(os.path.join(output_dir, fileName))
    return fileName
