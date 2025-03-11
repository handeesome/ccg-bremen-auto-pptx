import os
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from functions.helper import *
import re
from functions.今日事奉名单 import partOne
from functions.日期 import partTwo
from functions.主在圣殿中 import partThree
from functions.宣召 import partFour
from functions.敬拜赞美 import partFive
from functions.启应经文 import partSeven
from functions.主祷文 import partSix
from functions.读经 import partEight
from functions.讲道 import partNine
from functions.回应诗歌 import partTen
from functions.奉献回应礼 import partEleven
from functions.活动报告 import partTwelve
from functions.每月金句 import partThirteen
from functions.下周事奉名单 import partFourteen
from functions.欢迎新朋友 import partFifteen
from functions.圣餐礼 import partSixteen
from functions.三一颂 import partSeventeen
from functions.牧师祝福 import partEighteen
from functions.阿们颂 import partNineteen
from functions.默祷 import partTwenty
from functions.祷告会 import partTwentyOne


def generate_pptx(templatePath, formData):
    prs = Presentation(templatePath)
    # 今日事奉名单
    slide = prs.slides[0]
    证道 = ''.join(formData['thisWeekListzhengDao'])
    司会 = ''.join(formData['thisWeekListsiHui'])
    PPT = ''.join(formData['thisWeekListppt'])
    接待 = ''.join(formData['thisWeekListjieDai'])
    儿童主日学 = ''.join(formData['thisWeekListerTong'])
    partOne(slide, 证道, 司会, PPT, 接待, 儿童主日学)

    # 日期
    slide = prs.slides[1]
    date = formData['date1'].split('-')
    year = date[0]
    month = date[1]
    day = date[2]
    isCommunion = (formData["isCommunion"]=="圣餐崇拜")
    partTwo(slide, year, month, day, communion=isCommunion)

    # 主在圣殿中
    partThree(prs)

    # 宣召
    
    partFour(prs, formData['xuanZhao'])

    # 敬拜赞美 
    songOne = formData['song1']
    songTwo = formData['song2']
    songThree = formData['song3']
    songs = [
        formData.get('song1Pages', []),
        formData.get('song2Pages', []),
        formData.get('song3Pages', [])
    ]
    partFive(prs, songOne, songTwo, songThree, songs)

    # 为儿童祷告
    slide = newSlide(prs, '')
    slide.shapes.add_picture('docs/为儿童祷告.jpg', 0, 0, Cm(25.4), Cm(19.05))
    # 主祷文
    partSix(prs)

    # 启应经文
    qiyingLines = formData['qiYing']
    partSeven(prs, qiyingLines)

    # 读经 
    dujingLines = formData['duJing']
    partEight(prs, dujingLines)

    # 讲道
    title = formData['zhuTi1']
    经文 =''.join([jingWen["abbrName"] for jingWen in formData['jingWen']])
    partNine(prs, title, 证道, 经文)

    # 回应诗歌 
    songHuiying = formData['song4']
    songHuiyingPages = formData.get('song4Pages', [])
    partTen(prs, songHuiying, songHuiyingPages)

    # 奉献回应礼
    fengxianjingwen = formData['verseRadio']
    partEleven(prs, fengxianjingwen)

    # 活动报告
    activities = formData['activity']
    partTwelve(prs, activities, isCommunion=isCommunion)

    # 每月金句
    jinjuLines = formData['jinJu']
    partThirteen(prs, int(month), jinjuLines)

    # 下周事奉名单
    日期 = formData['date2']
    主题 = formData['zhuTi2']
    证道 = ' '.join(formData['nextWeekListzhengDao'])
    经文 = '\n'.join(verses["fullName"] for verses in formData['jingWen'])
    司会 = ' '.join(formData['nextWeekListsiHui'])
    PPT = ' '.join(formData['nextWeekListppt'])
    接待 = ' '.join(formData['nextWeekListjieDai'])
    儿童主日学 = ' '.join(formData['nextWeekListerTong'])
    partFourteen(prs, 日期, 主题, 证道, 经文, 司会, PPT, 接待, 儿童主日学)

    # 欢迎新朋友
    partFifteen(prs)

    # 圣餐礼
    shengcanSong = formData['lyricsRadio']
    if isCommunion:
        partSixteen(prs, shengcanSong)

    # 三一颂
    partSeventeen(prs)

    # 牧师祝福
    partEighteen(prs)

    # 阿们颂
    partNineteen(prs)

    # 默祷
    partTwenty(prs)

    # 祷告会
    prayerWorld = formData['prayerWorld']
    prayerChurch = formData['prayerChurch']
    partTwentyOne(prs, prayerWorld, prayerChurch)

    fileName = formData['date1'] + '.pptx'  # Name the file based on the date
    filePath = os.path.join("static/temp", fileName) 
    prs.save(filePath)  # Save the file
    print(filePath + " 生成成功，请在 temp/ 目录中查看")

    return fileName  # Return the filename for Flask to use

